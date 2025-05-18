from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import logging
import os
import glob
import traceback
from typing import Optional, List
from ..models import Presentation, Slide

# Configure logging
log_dir = "logs"
if not os.path.exists(log_dir):
    os.makedirs(log_dir)

logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(os.path.join(log_dir, 'presentation_generator.log')),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class PresentationGenerator:
    def __init__(self):
        try:
            self.prs = PPTXPresentation()
            # Set 16:9 aspect ratio
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
            
            self.colors = {
                'primary': RGBColor(0x00, 0x72, 0xC6),    # Microsoft Blue
                'secondary': RGBColor(0x40, 0x40, 0x40),   # Dark Gray
                'accent': RGBColor(0x00, 0xB2, 0x94),     # Teal
                'background': RGBColor(0x01, 0x16, 0x40),  # Dark Navy
                'text': RGBColor(0xFF, 0xFF, 0xFF)        # White
            }
            
            # Create presentations directory if it doesn't exist
            self.output_dir = "presentations"
            os.makedirs(self.output_dir, exist_ok=True)
            logger.info("PresentationGenerator initialized successfully")
            
            # Clean up old presentations
            self._cleanup_old_presentations()
        except Exception as e:
            logger.error(f"Error initializing PresentationGenerator: {str(e)}")
            logger.error(f"Stack trace: {traceback.format_exc()}")
            raise

    def generate(self, presentation: Presentation) -> Optional[str]:
        """Generate a PowerPoint presentation."""
        try:
            logger.info(f"Starting presentation generation: {presentation.title}")
            
            # Create a new presentation for each generation
            self.prs = PPTXPresentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
            
            # Process each slide
            for i, slide in enumerate(presentation.slides):
                logger.debug(f"Processing slide {i+1}: {slide.title}")
                try:
                    if i == 0 or slide.type == 'title':
                        self._create_title_slide(slide)
                    elif slide.type == 'table':
                        self._create_table_slide(slide)
                    else:
                        self._create_content_slide(slide)
                except Exception as e:
                    logger.error(f"Error creating slide {i+1}: {str(e)}")
                    logger.error(f"Stack trace: {traceback.format_exc()}")
                    # Continue with next slide instead of failing entire presentation
                    continue
            
            # Generate unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}.pptx"
            filepath = os.path.join(self.output_dir, filename)
            
            # Save the presentation
            try:
                self.prs.save(filepath)
                logger.info(f"Presentation saved successfully: {filepath}")
                return filename
            except Exception as e:
                logger.error(f"Error saving presentation: {str(e)}")
                logger.error(f"Stack trace: {traceback.format_exc()}")
                return None
                
        except Exception as e:
            logger.error(f"Error generating presentation: {str(e)}")
            logger.error(f"Stack trace: {traceback.format_exc()}")
            return None

    def _cleanup_old_presentations(self):
        """Clean up old presentation files."""
        try:
            logger.info("Starting presentation cleanup")
            # Clean up files in root directory
            root_files = glob.glob("presentation_*.pptx")
            for file in root_files:
                try:
                    os.remove(file)
                    logger.info(f"Cleaned up old presentation: {file}")
                except Exception as e:
                    logger.error(f"Error cleaning up file {file}: {str(e)}")
                    logger.error(f"Stack trace: {traceback.format_exc()}")
            
            # Clean up files in presentations directory
            pres_files = glob.glob(os.path.join(self.output_dir, "presentation_*.pptx"))
            # Keep only the 5 most recent files
            if len(pres_files) > 5:
                pres_files.sort(key=os.path.getctime)
                for file in pres_files[:-5]:  # Keep only last 5 files
                    try:
                        os.remove(file)
                        logger.info(f"Cleaned up old presentation: {file}")
                    except Exception as e:
                        logger.error(f"Error cleaning up file {file}: {str(e)}")
                        logger.error(f"Stack trace: {traceback.format_exc()}")
        except Exception as e:
            logger.error(f"Error during presentation cleanup: {str(e)}")
            logger.error(f"Stack trace: {traceback.format_exc()}")

    def _add_background_style(self, slide):
        """Add modern background style to slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']

        # Add decorative shapes - positioned better for visual appeal
        shape_width = Inches(5)
        shape_height = Inches(5)
        
        # Add network-like decorative elements
        positions = [
            (-2, -1, 15),  # Top left
            (11, -2, -15),  # Top right
            (9, 5, 30),    # Bottom right
        ]
        
        for x, y, rotation in positions:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x),
                Inches(y),
                shape_width,
                shape_height
            )
            shape_fill = shape.fill
            shape_fill.solid()
            shape_fill.fore_color.rgb = self.colors['primary']
            shape_fill.transparency = 0.85
            shape.rotation = rotation

    def _create_title_slide(self, slide: Slide) -> None:
        """Create a title slide with modern styling."""
        pptx_slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self._add_background_style(pptx_slide)
        
        try:
            # Center the title and subtitle in the middle of the slide
            title_left = Inches(1.5)
            title_top = Inches(2.0)  # Moved up slightly
            title_width = Inches(10.333)
            title_height = Inches(1.5)

            # Add title shape
            title_shape = pptx_slide.shapes.add_textbox(
                title_left, title_top, title_width, title_height
            )
            title_frame = title_shape.text_frame
            title_frame.clear()
            title_frame.word_wrap = True
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
            
            # Style title
            p = title_frame.add_paragraph()
            p.text = slide.title
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = self.colors['text']
            p.font.name = 'Segoe UI Light'
            
            # Calculate subtitle position based on title height
            subtitle_top = title_top + Inches(2.0)  # Dynamic spacing from title
            subtitle_height = Inches(1.2)
            
            # Add subtitle
            subtitle_shape = pptx_slide.shapes.add_textbox(
                title_left, subtitle_top, title_width, subtitle_height
            )
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.clear()
            subtitle_frame.word_wrap = True
            subtitle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            if slide.content:
                p = subtitle_frame.add_paragraph()
                content_text = "\n".join(slide.content) if isinstance(slide.content, list) else str(slide.content)
                p.text = content_text
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(32)
                p.font.color.rgb = self.colors['accent']
                p.font.name = 'Segoe UI'
                p.space_after = Pt(32)

                # Calculate line position based on subtitle content
                line_top = subtitle_top + subtitle_height + Inches(0.5)
            else:
                # If no subtitle content, position line closer to title
                line_top = subtitle_top + Inches(0.5)

            # Add decorative line with dynamic positioning
            line = pptx_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4),
                line_top,
                Inches(5.333),
                Inches(0.05)
            )
            line_fill = line.fill
            line_fill.solid()
            line_fill.fore_color.rgb = self.colors['accent']
            line.line.fill.background()

        except Exception as e:
            logger.error(f"Error creating title slide: {e}")

    def _create_content_slide(self, slide: Slide) -> None:
        """Create a content slide with modern styling."""
        pptx_slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_background_style(pptx_slide)
        
        try:
            # Define content area dimensions with better margins
            left_margin = Inches(1.8)
            content_width = Inches(9.733)
            
            # Add title with better positioning
            title_top = Inches(0.8)
            title_height = Inches(1.2)
            title_shape = pptx_slide.shapes.add_textbox(
                left_margin, title_top, content_width, title_height
            )
            title_frame = title_shape.text_frame
            title_frame.clear()
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = title_frame.add_paragraph()
            p.text = slide.title
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self.colors['text']
            p.font.name = 'Segoe UI'
            p.space_after = Pt(32)

            # Add content with better spacing and alignment
            if slide.content:
                content_top = title_top + Inches(1.4)
                content_height = Inches(5.0)
                content_shape = pptx_slide.shapes.add_textbox(
                    left_margin + Inches(0.5), content_top, 
                    content_width - Inches(1.0), content_height
                )
                text_frame = content_shape.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.TOP

                for point in slide.content:
                    p = text_frame.add_paragraph()
                    p.text = str(point)
                    p.font.size = Pt(24)
                    p.font.color.rgb = self.colors['text']
                    p.font.name = 'Segoe UI'
                    p.level = 0
                    p.space_after = Pt(16)  # Increased spacing between points
                    p.space_before = Pt(8)  # Added spacing before points
                    # Add bullet points
                    p.bullet.size = Pt(12)
                    p.bullet.font.color.rgb = self.colors['accent']

            # Add accent line
            accent = pptx_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8),
                Inches(0.8),
                Inches(0.1),
                Inches(5.9)
            )
            accent_fill = accent.fill
            accent_fill.solid()
            accent_fill.fore_color.rgb = self.colors['accent']
            accent.line.fill.background()

        except Exception as e:
            logger.error(f"Error creating content slide: {str(e)}")
            logger.error(f"Stack trace: {traceback.format_exc()}")

    def _create_table_slide(self, slide: Slide) -> None:
        """Create a slide with a table."""
        pptx_slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Using blank layout
        self._add_background_style(pptx_slide)
        
        try:
            # Add title with proper spacing
            title_left = Inches(1.5)
            title_top = Inches(0.8)
            title_width = Inches(10.333)
            title_height = Inches(1.0)
            
            title_shape = pptx_slide.shapes.add_textbox(
                title_left, title_top, title_width, title_height
            )
            title_frame = title_shape.text_frame
            title_frame.clear()
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = title_frame.add_paragraph()
            p.text = slide.title
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = self.colors['text']
            p.font.name = 'Segoe UI'
            p.space_after = Pt(24)

            # Process table content
            if isinstance(slide.content, list):
                # Calculate table dimensions
                rows = slide.content
                num_rows = len(rows)
                num_cols = len(rows[0]) if rows else 2
                
                # Create table with better positioning
                table_left = Inches(1.5)
                table_top = Inches(2.0)
                table_width = Inches(10.333)
                row_height = Inches(0.6)
                
                table = pptx_slide.shapes.add_table(
                    num_rows, num_cols,
                    table_left, table_top,
                    table_width, row_height * num_rows
                ).table

                # Set column widths evenly
                col_width = table_width / num_cols
                for col in table.columns:
                    col.width = col_width
                
                # Add data to table with improved formatting
                for row_idx, row_data in enumerate(rows):
                    for col_idx, cell_data in enumerate(row_data):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_data)
                        
                        # Apply cell formatting
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = self.colors['text']
                        paragraph.font.name = 'Segoe UI'
                        
                        # Center align headers (first row)
                        if row_idx == 0:
                            paragraph.alignment = PP_ALIGN.CENTER
                            paragraph.font.bold = True
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = self.colors['primary']
                        else:
                            paragraph.alignment = PP_ALIGN.LEFT
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = self.colors['background'] if row_idx % 2 == 0 else self.colors['secondary']
                        
                        # Adjust cell margins and vertical alignment
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                        cell.margin_left = Inches(0.1)
                        cell.margin_right = Inches(0.1)
                        cell.margin_top = Inches(0.05)
                        cell.margin_bottom = Inches(0.05)

            # Add accent line
            accent = pptx_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8),
                Inches(0.8),
                Inches(0.1),
                Inches(5.9)
            )
            accent_fill = accent.fill
            accent_fill.solid()
            accent_fill.fore_color.rgb = self.colors['accent']
            accent.line.fill.background()

        except Exception as e:
            logger.error(f"Error creating table slide: {str(e)}")
            logger.error(f"Stack trace: {traceback.format_exc()}") 