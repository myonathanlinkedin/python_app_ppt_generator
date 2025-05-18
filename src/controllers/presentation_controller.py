from flask import Blueprint, jsonify, request, send_file, session, render_template
import logging
import os
import glob
import json
from typing import Tuple, Union
from datetime import datetime, timedelta
from weasyprint import HTML, CSS
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from ..services.llm_service import LLMService
from ..services.presentation_generator import PresentationGenerator
from ..models import Presentation, Theme, Slide
import io

# Configure logging to write to both file and console
log_dir = "logs"
if not os.path.exists(log_dir):
    os.makedirs(log_dir)

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(os.path.join(log_dir, 'app.log')),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class PresentationController:
    def __init__(self):
        self.blueprint = Blueprint('presentation', __name__)
        self.llm_service = LLMService()
        self.presentation_generator = PresentationGenerator()
        self.output_dir = "presentations"
        
        # Create output directory if it doesn't exist
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)

        # Register routes
        self.blueprint.route('/generate', methods=['POST'])(self.generate_preview)
        self.blueprint.route('/export/pdf', methods=['POST'])(self.export_pdf)
        self.blueprint.route('/export/ppt', methods=['POST'])(self.export_pptx)
        
        # Clean up old presentations
        self.cleanup_old_presentations()
        logger.info("PresentationController initialized")

    def cleanup_old_presentations(self):
        """Remove presentations older than 1 hour."""
        try:
            current_time = datetime.now()
            pattern = os.path.join(self.output_dir, "*.pdf")
            for file_path in glob.glob(pattern):
                creation_time = datetime.fromtimestamp(os.path.getctime(file_path))
                if current_time - creation_time > timedelta(hours=1):
                    try:
                        os.remove(file_path)
                        logger.info(f"Removed old presentation: {file_path}")
                    except Exception as e:
                        logger.error(f"Error removing old presentation {file_path}: {str(e)}")
        except Exception as e:
            logger.error(f"Error during presentation cleanup: {str(e)}")

    def generate_preview(self):
        """Generate presentation content and return HTML preview."""
        try:
            data = request.get_json()
            if not data:
                logger.error("No data provided in request")
                return jsonify({'error': 'No data provided'}), 400
                
            topic = data.get('topic')
            if not topic or not topic.strip():
                logger.error("No topic provided")
                return jsonify({'error': 'Please provide a presentation topic'}), 400
                
            style = data.get('style', 'corporate')
            logger.info(f"Generating presentation preview for topic: {topic}, style: {style}")
            
            try:
                # Generate content using LLM service
                content_data = self.llm_service.generate_presentation_content(topic)
                
                # Add theme if not present
                if 'theme' not in content_data:
                    content_data['theme'] = {
                        'primary_color': '#0072C6',
                        'secondary_color': '#404040',
                        'accent_color': '#00B294',
                        'background_color': '#FFFFFF'
                    }

                return jsonify({
                    'presentation': content_data
                })
                
            except ValueError as e:
                # Handle validation errors with specific messages
                error_msg = str(e)
                logger.error(f"Validation error: {error_msg}")
                return jsonify({
                    'error': error_msg
                }), 400
                
            except ConnectionError as e:
                # Handle LLM service connection errors
                error_msg = str(e)
                logger.error(f"LLM service error: {error_msg}")
                return jsonify({
                    'error': error_msg
                }), 503
                
            except Exception as e:
                # Handle unexpected errors
                logger.error(f"Unexpected error in LLM service: {str(e)}", exc_info=True)
                return jsonify({
                    'error': 'An unexpected error occurred while generating the presentation. Please try again with a different topic.'
                }), 500

        except Exception as e:
            # Handle request processing errors
            logger.error(f"Error processing request: {str(e)}", exc_info=True)
            return jsonify({
                'error': 'An error occurred while processing your request. Please try again.'
            }), 500

    def export_pdf(self):
        """Export presentation as PDF."""
        try:
            data = request.get_json()
            if not data or 'presentation' not in data:
                return jsonify({'error': 'No presentation data provided'}), 400

            # Generate HTML with print-optimized styles
            html = render_template('presentation.html', 
                                 presentation=data['presentation'],
                                 print_mode=True)

            # Create temporary HTML file for WeasyPrint
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            temp_html_path = os.path.join(self.output_dir, f"temp_{timestamp}.html")
            
            with open(temp_html_path, 'w', encoding='utf-8') as f:
                f.write(html)

            # Configure WeasyPrint with custom settings
            css = CSS(string='''
                @page {
                    size: 1280px 720px;
                    margin: 0;
                }
                body { margin: 0; }
            ''')
            
            pdf = HTML(filename=temp_html_path).write_pdf(
                stylesheets=[css],
                presentational_hints=True
            )
            
            # Clean up temporary HTML file
            os.remove(temp_html_path)
            
            # Return PDF as attachment
            return send_file(
                io.BytesIO(pdf),
                mimetype='application/pdf',
                as_attachment=True,
                download_name='presentation.pdf'
            )

        except Exception as e:
            logger.error(f"Error in export_pdf: {str(e)}", exc_info=True)
            return jsonify({'error': f'Internal server error: {str(e)}'}), 500

    def export_pptx(self):
        """Export presentation as PowerPoint."""
        try:
            data = request.get_json()
            if not data or 'presentation' not in data:
                return jsonify({'error': 'No presentation data provided'}), 400

            # Create PowerPoint presentation
            prs = PPTXPresentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Set theme colors
            theme = data['presentation'].get('theme', {})
            colors = {
                'primary': RGBColor.from_string(theme.get('primary_color', '#0072C6').lstrip('#')),
                'secondary': RGBColor.from_string(theme.get('secondary_color', '#404040').lstrip('#')),
                'accent': RGBColor.from_string(theme.get('accent_color', '#00B294').lstrip('#')),
                'background': RGBColor.from_string(theme.get('background_color', '#FFFFFF').lstrip('#')),
                'text': RGBColor(0x00, 0x00, 0x00)  # Black
            }

            # Create slides
            for slide_data in data['presentation']['slides']:
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                shapes = slide.shapes
                
                # Add title
                title = shapes.title
                title.text = slide_data['title']
                title.text_frame.paragraphs[0].font.size = Pt(36)
                title.text_frame.paragraphs[0].font.color.rgb = colors['primary']
                
                # Add content based on slide type
                if slide_data['type'] == 'table':
                    # Add table
                    rows = len(slide_data['content'])
                    cols = len(slide_data['content'][0]) if rows > 0 else 2
                    
                    left = Inches(2.0)
                    top = Inches(2.0)
                    width = Inches(9.0)
                    height = Inches(4.0)
                    
                    table = shapes.add_table(rows, cols, left, top, width, height).table
                    
                    # Populate table
                    for row_idx, row in enumerate(slide_data['content']):
                        for col_idx, cell_text in enumerate(row):
                            cell = table.cell(row_idx, col_idx)
                            cell.text = str(cell_text)
                            cell.text_frame.paragraphs[0].font.size = Pt(14)
                            cell.text_frame.paragraphs[0].font.color.rgb = colors['text']
                else:
                    # Add text content
                    content = slide_data.get('content', [])
                    if isinstance(content, list):
                        text_box = shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
                        text_frame = text_box.text_frame
                        
                        for point in content:
                            p = text_frame.add_paragraph()
                            p.text = str(point)
                            p.font.size = Pt(18)
                            p.font.color.rgb = colors['text']
                            p.level = 0
            
            # Save to BytesIO
            pptx_stream = io.BytesIO()
            prs.save(pptx_stream)
            pptx_stream.seek(0)
            
            # Return PPTX as attachment
            return send_file(
                pptx_stream,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name='presentation.pptx'
            )

        except Exception as e:
            logger.error(f"Error in export_pptx: {str(e)}", exc_info=True)
            return jsonify({'error': f'Internal server error: {str(e)}'}), 500 